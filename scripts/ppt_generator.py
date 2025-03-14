# ppt_generator.py
from pptx import Presentation
import random
import time
from pathlib import Path

def generate_ppt_all():
    types = ['商务', '教育']
    
    for template_type in types:
        output_dir = Path(f"D:/AI_System/data/raw/电子模板素材/可编辑PPT/{template_type}")
        output_dir.mkdir(parents=True, exist_ok=True)
        
        for i in range(5):  # 每个类型生成5个
            try:
                prs = Presentation()
                prs.slide_width = 12192000  # 16:9宽屏
                prs.slide_height = 6858000
                
                # 智能选择带标题的版式
                valid_layouts = [lyt for lyt in prs.slide_layouts if lyt.name in ['Title Slide', 'Title and Content']]
                
                for _ in range(random.randint(3,6)):
                    slide_layout = random.choice(valid_layouts)
                    slide = prs.slides.add_slide(slide_layout)
                    
                    # 确保标题对象存在
                    if slide.shapes.title:
                        title = slide.shapes.title
                        title.text = f"{template_type}模板-{i+1}"
                    else:
                        title_box = slide.shapes.add_textbox(left=0, top=0, width=12192000, height=685800)
                        title = title_box.text_frame.add_paragraph()
                        title.text = f"{template_type}模板-{i+1}"
                
                # 唯一文件名（毫秒级时间戳）
                timestamp = int(time.time() * 1000) + i
                prs.save(output_dir / f"{template_type}_模板_{timestamp}.pptx")
                print(f"✅ {template_type}模板已生成：{template_type}_模板_{timestamp}.pptx")
                
            except Exception as e:
                print(f"❌ {template_type}模板生成失败：{str(e)}")

if __name__ == "__main__":
    generate_ppt_all()